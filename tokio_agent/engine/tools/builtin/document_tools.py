"""
Document Tools — Generate PDF reports, PPTX slides, and CSV exports.

Uses ReportLab for PDF (professional tables, colors, borders),
python-pptx for slides, csv module for CSV.
Generated files are saved to /workspace/output/ for Telegram delivery.
"""
from __future__ import annotations

import csv
import io
import json
import logging
import os
import re
from datetime import datetime
from pathlib import Path
from typing import Any, Dict, List, Optional

logger = logging.getLogger(__name__)

OUTPUT_DIR = Path(os.getenv("TOKIO_OUTPUT_DIR", "/workspace/output"))


def _ensure_output_dir() -> Path:
    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
    return OUTPUT_DIR


# ── Text Processing ──────────────────────────────────────────────────────

def _sanitize_text(text: str) -> str:
    """Replace problematic Unicode chars with ASCII equivalents."""
    replacements = {
        "\u2014": "-", "\u2013": "-",
        "\u2018": "'", "\u2019": "'",
        "\u201c": '"', "\u201d": '"',
        "\u2026": "...", "\u2022": "*",
        "\u00b7": "*", "\u2192": "->",
        "\u2190": "<-", "\u2264": "<=",
        "\u2265": ">=", "\u2260": "!=",
    }
    for old, new in replacements.items():
        text = text.replace(old, new)
    return text


def _parse_markdown_table(text: str) -> Optional[List[List[str]]]:
    """Extract a markdown table from text. Returns list of rows including header."""
    lines = text.strip().split("\n")
    table_lines = []
    in_table = False
    for line in lines:
        stripped = line.strip()
        if "|" in stripped and stripped.startswith("|"):
            # Skip separator lines like |---|---|
            if re.match(r'^\|[\s\-:|]+\|$', stripped):
                in_table = True
                continue
            cells = [c.strip() for c in stripped.split("|")[1:-1]]
            if cells:
                table_lines.append(cells)
                in_table = True
        elif in_table and not stripped:
            break
    return table_lines if table_lines else None


def _extract_tables_and_text(body: str) -> list:
    """Parse body text and extract tables vs regular text blocks.
    Returns list of ('text', str) or ('table', [[rows]]) tuples.
    """
    lines = body.split("\n")
    blocks = []
    current_text = []
    current_table = []
    in_table = False
    
    for line in lines:
        stripped = line.strip()
        is_table_line = "|" in stripped and stripped.startswith("|") and stripped.endswith("|")
        is_separator = bool(re.match(r'^\|[\s\-:|]+\|$', stripped))
        
        if is_table_line or is_separator:
            # Flush text
            if current_text:
                text = "\n".join(current_text).strip()
                if text:
                    blocks.append(("text", text))
                current_text = []
            
            if is_separator:
                in_table = True
                continue
            
            cells = [c.strip() for c in stripped.split("|")[1:-1]]
            if cells:
                current_table.append(cells)
            in_table = True
        else:
            if in_table and current_table:
                blocks.append(("table", current_table))
                current_table = []
                in_table = False
            current_text.append(line)
    
    # Flush remaining
    if current_table:
        blocks.append(("table", current_table))
    if current_text:
        text = "\n".join(current_text).strip()
        if text:
            blocks.append(("text", text))
    
    return blocks


# ── PDF Generation (ReportLab) ───────────────────────────────────────────

def _generate_pdf(
    title: str,
    sections: List[Dict[str, str]],
    output_path: str = "",
    template: str = "default",
) -> str:
    """Generate a professional PDF report using ReportLab."""
    try:
        from reportlab.lib import colors
        from reportlab.lib.pagesizes import A4
        from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
        from reportlab.lib.units import mm
        from reportlab.platypus import (
            SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle,
            PageBreak, HRFlowable
        )
        from reportlab.lib.enums import TA_CENTER, TA_LEFT
    except ImportError:
        return json.dumps({"ok": False, "error": "reportlab no instalado. Ejecuta: pip install reportlab"})

    if not output_path:
        safe_title = "".join(c if c.isalnum() or c in " _-" else "_" for c in title)[:60]
        output_path = str(_ensure_output_dir() / f"{safe_title}_{datetime.now().strftime('%Y%m%d_%H%M')}.pdf")

    # Color schemes
    schemes = {
        "default": {
            "title_bg": colors.HexColor("#1a1a2e"),
            "title_fg": colors.white,
            "heading_bg": colors.HexColor("#2980b9"),
            "heading_fg": colors.white,
            "table_header_bg": colors.HexColor("#2c3e50"),
            "table_header_fg": colors.white,
            "table_alt_bg": colors.HexColor("#ecf0f1"),
            "accent": colors.HexColor("#3498db"),
        },
        "security": {
            "title_bg": colors.HexColor("#1a0000"),
            "title_fg": colors.white,
            "heading_bg": colors.HexColor("#c0392b"),
            "heading_fg": colors.white,
            "table_header_bg": colors.HexColor("#7b241c"),
            "table_header_fg": colors.white,
            "table_alt_bg": colors.HexColor("#f9ebea"),
            "accent": colors.HexColor("#e74c3c"),
        },
        "infrastructure": {
            "title_bg": colors.HexColor("#0a2e1a"),
            "title_fg": colors.white,
            "heading_bg": colors.HexColor("#27ae60"),
            "heading_fg": colors.white,
            "table_header_bg": colors.HexColor("#1e8449"),
            "table_header_fg": colors.white,
            "table_alt_bg": colors.HexColor("#eafaf1"),
            "accent": colors.HexColor("#2ecc71"),
        },
    }
    scheme = schemes.get(template, schemes["default"])

    # Build document
    Path(output_path).parent.mkdir(parents=True, exist_ok=True)
    doc = SimpleDocTemplate(
        output_path,
        pagesize=A4,
        leftMargin=20*mm,
        rightMargin=20*mm,
        topMargin=15*mm,
        bottomMargin=15*mm,
    )

    styles = getSampleStyleSheet()
    
    # Custom styles
    styles.add(ParagraphStyle(
        name="TitleBanner",
        parent=styles["Title"],
        fontSize=22,
        textColor=scheme["title_fg"],
        alignment=TA_CENTER,
        spaceAfter=6,
    ))
    styles.add(ParagraphStyle(
        name="Subtitle",
        parent=styles["Normal"],
        fontSize=10,
        textColor=colors.HexColor("#cccccc"),
        alignment=TA_CENTER,
        spaceAfter=4,
    ))
    styles.add(ParagraphStyle(
        name="SectionHeading",
        parent=styles["Heading2"],
        fontSize=13,
        textColor=scheme["heading_fg"],
        spaceAfter=4,
        spaceBefore=8,
        leftIndent=4,
    ))
    styles.add(ParagraphStyle(
        name="BodyText2",
        parent=styles["Normal"],
        fontSize=10,
        leading=14,
        textColor=colors.HexColor("#333333"),
        spaceAfter=4,
    ))
    styles.add(ParagraphStyle(
        name="BulletItem",
        parent=styles["Normal"],
        fontSize=10,
        leading=14,
        textColor=colors.HexColor("#333333"),
        leftIndent=15,
        spaceAfter=2,
    ))
    styles.add(ParagraphStyle(
        name="FooterText",
        parent=styles["Normal"],
        fontSize=8,
        textColor=colors.HexColor("#999999"),
        alignment=TA_CENTER,
    ))
    styles.add(ParagraphStyle(
        name="TableCell",
        parent=styles["Normal"],
        fontSize=9,
        leading=12,
        textColor=colors.HexColor("#333333"),
    ))
    styles.add(ParagraphStyle(
        name="TableHeaderCell",
        parent=styles["Normal"],
        fontSize=9,
        leading=12,
        textColor=colors.white,
        fontName="Helvetica-Bold",
    ))

    elements = []
    page_width = A4[0] - 40*mm  # Available width

    # ── Title Banner ──
    title_data = [[Paragraph(title, styles["TitleBanner"])]]
    title_table = Table(title_data, colWidths=[page_width])
    title_table.setStyle(TableStyle([
        ("BACKGROUND", (0, 0), (-1, -1), scheme["title_bg"]),
        ("TOPPADDING", (0, 0), (-1, -1), 16),
        ("BOTTOMPADDING", (0, 0), (-1, -1), 8),
        ("LEFTPADDING", (0, 0), (-1, -1), 10),
        ("RIGHTPADDING", (0, 0), (-1, -1), 10),
        ("ROUNDEDCORNERS", [4, 4, 4, 4]),
    ]))
    elements.append(title_table)

    # Date and source
    sub_data = [[Paragraph(
        f"Generado: {datetime.now().strftime('%Y-%m-%d %H:%M')}  |  TokioAI v2.1",
        styles["Subtitle"]
    )]]
    sub_table = Table(sub_data, colWidths=[page_width])
    sub_table.setStyle(TableStyle([
        ("BACKGROUND", (0, 0), (-1, -1), scheme["title_bg"]),
        ("TOPPADDING", (0, 0), (-1, -1), 0),
        ("BOTTOMPADDING", (0, 0), (-1, -1), 10),
        ("LEFTPADDING", (0, 0), (-1, -1), 10),
        ("RIGHTPADDING", (0, 0), (-1, -1), 10),
        ("ROUNDEDCORNERS", [0, 0, 4, 4]),
    ]))
    elements.append(sub_table)
    elements.append(Spacer(1, 12))

    # ── Sections ──
    for section in sections:
        heading = section.get("heading", "")
        body = section.get("body", "")

        if heading:
            h_data = [[Paragraph(f"  {heading}", styles["SectionHeading"])]]
            h_table = Table(h_data, colWidths=[page_width])
            h_table.setStyle(TableStyle([
                ("BACKGROUND", (0, 0), (-1, -1), scheme["heading_bg"]),
                ("TOPPADDING", (0, 0), (-1, -1), 6),
                ("BOTTOMPADDING", (0, 0), (-1, -1), 6),
                ("LEFTPADDING", (0, 0), (-1, -1), 6),
                ("ROUNDEDCORNERS", [3, 3, 3, 3]),
            ]))
            elements.append(h_table)
            elements.append(Spacer(1, 6))

        if body:
            blocks = _extract_tables_and_text(body)
            for block_type, block_content in blocks:
                if block_type == "table" and len(block_content) >= 1:
                    _add_table(elements, block_content, scheme, styles, page_width)
                elif block_type == "text":
                    for line in block_content.split("\n"):
                        line = line.strip()
                        if not line:
                            elements.append(Spacer(1, 4))
                        elif line.startswith(("- ", "* ", "-> ", "=> ")):
                            bullet_text = line.lstrip("-*> ").strip()
                            elements.append(Paragraph(
                                f"\u2022  {_escape_html(bullet_text)}", 
                                styles["BulletItem"]
                            ))
                        elif line.startswith(("#", "##", "###")):
                            clean = line.lstrip("# ").strip()
                            elements.append(Paragraph(
                                f"<b>{_escape_html(clean)}</b>", 
                                styles["BodyText2"]
                            ))
                        else:
                            elements.append(Paragraph(
                                _escape_html(line), 
                                styles["BodyText2"]
                            ))

            elements.append(Spacer(1, 8))

    # ── Footer ──
    elements.append(Spacer(1, 20))
    elements.append(HRFlowable(width="100%", thickness=1, color=colors.HexColor("#cccccc")))
    elements.append(Spacer(1, 4))
    elements.append(Paragraph(
        "Generado por TokioAI - Agente Autonomo de Seguridad",
        styles["FooterText"]
    ))

    # Build PDF
    try:
        doc.build(elements)
        size = os.path.getsize(output_path)
        return json.dumps({
            "ok": True,
            "file": output_path,
            "size_bytes": size,
            "pages": "auto",
        }, ensure_ascii=False)
    except Exception as e:
        logger.error(f"Error building PDF: {e}")
        return json.dumps({"ok": False, "error": f"Error generando PDF: {e}"})


def _escape_html(text: str) -> str:
    """Escape HTML special chars for ReportLab Paragraph."""
    text = text.replace("&", "&amp;")
    text = text.replace("<", "&lt;")
    text = text.replace(">", "&gt;")
    # Convert **bold** to <b>bold</b>
    text = re.sub(r'\*\*(.+?)\*\*', r'<b>\1</b>', text)
    # Convert `code` to <font color="#c0392b">code</font>
    text = re.sub(r'`(.+?)`', r'<font color="#c0392b" face="Courier">\1</font>', text)
    return text


def _add_table(elements, rows, scheme, styles, page_width):
    """Add a formatted table to the elements list."""
    from reportlab.lib import colors as rl_colors
    from reportlab.platypus import Table, TableStyle, Paragraph, Spacer
    
    if not rows:
        return
    
    # Determine columns
    max_cols = max(len(row) for row in rows)
    col_width = page_width / max_cols
    col_widths = [col_width] * max_cols
    
    # Build table data with Paragraphs
    table_data = []
    for i, row in enumerate(rows):
        # Pad row if needed
        while len(row) < max_cols:
            row.append("")
        
        if i == 0:
            # Header row
            table_data.append([
                Paragraph(_escape_html(str(cell)), styles["TableHeaderCell"])
                for cell in row
            ])
        else:
            table_data.append([
                Paragraph(_escape_html(str(cell)), styles["TableCell"])
                for cell in row
            ])
    
    table = Table(table_data, colWidths=col_widths, repeatRows=1)
    
    # Table style
    style_commands = [
        # Header
        ("BACKGROUND", (0, 0), (-1, 0), scheme["table_header_bg"]),
        ("TEXTCOLOR", (0, 0), (-1, 0), scheme["table_header_fg"]),
        ("FONTNAME", (0, 0), (-1, 0), "Helvetica-Bold"),
        ("FONTSIZE", (0, 0), (-1, 0), 9),
        # All cells
        ("FONTSIZE", (0, 1), (-1, -1), 9),
        ("TOPPADDING", (0, 0), (-1, -1), 5),
        ("BOTTOMPADDING", (0, 0), (-1, -1), 5),
        ("LEFTPADDING", (0, 0), (-1, -1), 6),
        ("RIGHTPADDING", (0, 0), (-1, -1), 6),
        # Grid
        ("GRID", (0, 0), (-1, -1), 0.5, rl_colors.HexColor("#bdc3c7")),
        ("LINEBELOW", (0, 0), (-1, 0), 1.5, scheme["heading_bg"]),
        # Alignment
        ("VALIGN", (0, 0), (-1, -1), "MIDDLE"),
    ]
    
    # Alternating row colors
    for i in range(1, len(table_data)):
        if i % 2 == 0:
            style_commands.append(
                ("BACKGROUND", (0, i), (-1, i), scheme["table_alt_bg"])
            )
    
    table.setStyle(TableStyle(style_commands))
    elements.append(table)
    elements.append(Spacer(1, 8))


# ── Slides Generation (python-pptx) ──────────────────────────────────────

def _generate_slides(
    title: str,
    slides: List[Dict[str, Any]],
    output_path: str = "",
    template: str = "default",
) -> str:
    """Generate PPTX slides with professional formatting."""
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt, Emu
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    except ImportError:
        return json.dumps({"ok": False, "error": "python-pptx no instalado. Ejecuta: pip install python-pptx"})

    if not output_path:
        safe_title = "".join(c if c.isalnum() or c in " _-" else "_" for c in title)[:60]
        output_path = str(_ensure_output_dir() / f"{safe_title}_{datetime.now().strftime('%Y%m%d_%H%M')}.pptx")

    # Color schemes
    schemes = {
        "default": {
            "bg": RGBColor(0x1a, 0x1a, 0x2e),
            "title_color": RGBColor(0xff, 0xff, 0xff),
            "accent": RGBColor(0x34, 0x98, 0xdb),
            "text": RGBColor(0xec, 0xf0, 0xf1),
            "bullet_color": RGBColor(0x2e, 0xcc, 0x71),
        },
        "security": {
            "bg": RGBColor(0x1a, 0x00, 0x00),
            "title_color": RGBColor(0xff, 0xff, 0xff),
            "accent": RGBColor(0xe7, 0x4c, 0x3c),
            "text": RGBColor(0xec, 0xf0, 0xf1),
            "bullet_color": RGBColor(0xe7, 0x4c, 0x3c),
        },
        "infrastructure": {
            "bg": RGBColor(0x0a, 0x2e, 0x1a),
            "title_color": RGBColor(0xff, 0xff, 0xff),
            "accent": RGBColor(0x2e, 0xcc, 0x71),
            "text": RGBColor(0xec, 0xf0, 0xf1),
            "bullet_color": RGBColor(0x27, 0xae, 0x60),
        },
    }
    scheme = schemes.get(template, schemes["default"])

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    def set_slide_bg(slide, color):
        """Set solid background color for a slide."""
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = color

    def add_text_box(slide, left, top, width, height, text, font_size=18,
                     color=None, bold=False, alignment=PP_ALIGN.LEFT):
        """Add a text box to a slide."""
        from pptx.util import Inches, Pt
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.color.rgb = color or scheme["text"]
        p.font.bold = bold
        p.alignment = alignment
        return tf

    # ── Title Slide ──
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    set_slide_bg(slide, scheme["bg"])
    
    # Title
    add_text_box(slide, Inches(1), Inches(2), Inches(11.333), Inches(1.5),
                 title, font_size=40, color=scheme["title_color"], bold=True,
                 alignment=PP_ALIGN.CENTER)
    
    # Accent line
    from pptx.util import Inches, Pt, Emu
    shape = slide.shapes.add_shape(
        1, Inches(4), Inches(3.7), Inches(5.333), Inches(0.05)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = scheme["accent"]
    shape.line.fill.background()
    
    # Subtitle
    add_text_box(slide, Inches(1), Inches(4), Inches(11.333), Inches(1),
                 f"Generado: {datetime.now().strftime('%Y-%m-%d %H:%M')}  |  TokioAI v2.1",
                 font_size=16, color=scheme["text"], alignment=PP_ALIGN.CENTER)

    # ── Content Slides ──
    for slide_data in slides:
        slide_title = slide_data.get("title", "")
        content = slide_data.get("content", "")
        bullets = slide_data.get("bullets", [])

        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
        set_slide_bg(slide, scheme["bg"])
        
        # Title bar
        title_shape = slide.shapes.add_shape(
            1, Inches(0), Inches(0), Inches(13.333), Inches(1.2)
        )
        title_shape.fill.solid()
        title_shape.fill.fore_color.rgb = RGBColor(
            max(0, scheme["bg"][0] - 10),
            max(0, scheme["bg"][1] - 10),
            max(0, scheme["bg"][2] - 10),
        )
        title_shape.line.fill.background()
        
        add_text_box(slide, Inches(0.8), Inches(0.15), Inches(11.733), Inches(0.9),
                     slide_title, font_size=28, color=scheme["accent"], bold=True)

        if bullets:
            y = 1.6
            for bullet in bullets:
                bullet_str = str(bullet)
                # Clean bullet markers
                bullet_str = bullet_str.lstrip("-*> ").strip()
                
                tf = add_text_box(
                    slide, Inches(1.2), Inches(y), Inches(11), Inches(0.5),
                    f"\u2022  {bullet_str}", font_size=18, color=scheme["text"]
                )
                y += 0.55
                if y > 6.5:
                    break
        elif content:
            add_text_box(slide, Inches(0.8), Inches(1.6), Inches(11.733), Inches(5.4),
                         content, font_size=16, color=scheme["text"])

    # Save
    try:
        Path(output_path).parent.mkdir(parents=True, exist_ok=True)
        prs.save(output_path)
        size = os.path.getsize(output_path)
        return json.dumps({
            "ok": True,
            "file": output_path,
            "size_bytes": size,
            "slides_count": len(prs.slides),
        }, ensure_ascii=False)
    except Exception as e:
        return json.dumps({"ok": False, "error": f"Error guardando PPTX: {e}"})


# ── CSV Generation ────────────────────────────────────────────────────────

def _generate_csv(
    data: List[List[Any]],
    output_path: str = "",
    headers: Optional[List[str]] = None,
) -> str:
    """Generate CSV file."""
    if not output_path:
        output_path = str(_ensure_output_dir() / f"export_{datetime.now().strftime('%Y%m%d_%H%M')}.csv")

    try:
        Path(output_path).parent.mkdir(parents=True, exist_ok=True)
        with open(output_path, "w", newline="", encoding="utf-8") as f:
            writer = csv.writer(f)
            if headers:
                writer.writerow(headers)
            for row in data:
                writer.writerow(row)
        size = os.path.getsize(output_path)
        return json.dumps({
            "ok": True,
            "file": output_path,
            "size_bytes": size,
            "rows": len(data),
        }, ensure_ascii=False)
    except Exception as e:
        return json.dumps({"ok": False, "error": f"Error guardando CSV: {e}"})


# ── Unified Tool Entry Point ─────────────────────────────────────────────

async def document_tool(
    action: str,
    title: Optional[str] = None,
    content: Optional[str] = None,
    sections: Optional[List[Dict[str, str]]] = None,
    data: Optional[List[List[Any]]] = None,
    headers: Optional[List[str]] = None,
    output_path: Optional[str] = None,
    template: Optional[str] = None,
    params: Optional[Dict[str, Any]] = None,
) -> str:
    """
    Document generation tool.

    Actions:
      - generate_pdf: Create PDF report (uses ReportLab with professional tables)
      - generate_slides: Create PPTX presentation (dark theme)
      - generate_csv: Export CSV
    """
    # Support legacy nested params format
    if params and isinstance(params, dict):
        title = title or params.get("title")
        content = content or params.get("content") or params.get("body") or params.get("text")
        sections = sections or params.get("sections")
        data = data or params.get("data")
        headers = headers or params.get("headers")
        output_path = output_path or params.get("output_path")
        template = template or params.get("template")

    action = (action or "").strip().lower()
    title = str(title or "Reporte TokioAI").strip()
    output_path = str(output_path or "").strip()
    template = str(template or "default").strip()

    try:
        if action == "generate_pdf":
            pdf_sections = []

            # Normalize sections
            if sections:
                if isinstance(sections, str):
                    pdf_sections = [{"heading": "", "body": sections}]
                elif isinstance(sections, list):
                    for s in sections:
                        if isinstance(s, str):
                            pdf_sections.append({"heading": "", "body": s})
                        elif isinstance(s, dict):
                            pdf_sections.append(s)

            # Fallback: use content as single body
            if not pdf_sections and content:
                if isinstance(content, list):
                    pdf_sections = [{"heading": "", "body": "\n".join(str(x) for x in content)}]
                else:
                    pdf_sections = [{"heading": "", "body": str(content)}]

            if not pdf_sections:
                return json.dumps({"ok": False, "error": (
                    "Se necesita 'content' (texto) o 'sections' "
                    "(lista de {heading, body}) para generar el PDF."
                )})

            return _generate_pdf(title, pdf_sections, output_path, template)

        elif action == "generate_slides":
            slides_list = []

            if sections:
                if isinstance(sections, str):
                    slides_list = [{"title": "Contenido", "content": sections}]
                elif isinstance(sections, list):
                    for s in sections:
                        if isinstance(s, str):
                            slides_list.append({"title": "", "content": s})
                        elif isinstance(s, dict):
                            slide_entry = {
                                "title": s.get("heading", s.get("title", "")),
                            }
                            body = s.get("body", s.get("content", ""))
                            bullets_raw = s.get("bullets", [])
                            
                            # Auto-extract bullets from body
                            if not bullets_raw and body:
                                lines = body.strip().split("\n")
                                extracted = []
                                for line in lines:
                                    line = line.strip()
                                    if line and not line.startswith("|"):
                                        extracted.append(line)
                                if extracted:
                                    bullets_raw = extracted
                            
                            if bullets_raw:
                                slide_entry["bullets"] = bullets_raw
                            else:
                                slide_entry["content"] = body
                            slides_list.append(slide_entry)

            if not slides_list and content:
                lines = str(content).strip().split("\n")
                slides_list = [{"title": title, "bullets": [l for l in lines if l.strip()]}]

            if not slides_list:
                return json.dumps({"ok": False, "error": (
                    "Se necesita 'sections' (lista de {heading/title, body/bullets}) "
                    "para generar slides."
                )})

            return _generate_slides(title, slides_list, output_path, template)

        elif action == "generate_csv":
            csv_data = []
            csv_headers = headers

            if data:
                if isinstance(data, str):
                    csv_data = json.loads(data)
                elif isinstance(data, list):
                    csv_data = data

            if not csv_data:
                return json.dumps({"ok": False, "error": "Se necesita 'data' (lista de filas) para generar CSV."})

            return _generate_csv(csv_data, output_path, csv_headers)

        else:
            return json.dumps({"ok": False, "error": f"Accion desconocida: '{action}'. Usa: generate_pdf, generate_slides, generate_csv"})

    except Exception as e:
        logger.error(f"Document tool error: {e}", exc_info=True)
        return json.dumps({"ok": False, "error": str(e)})
